import os
import tempfile
from typing import List, Optional, Tuple
from datetime import datetime
import fitz  # PyMuPDF
import docx
import pandas as pd
from pptx import Presentation
import boto3
from fastapi import UploadFile
import uuid
from pathlib import Path

from app.models.db_models import FileType, Chunk
from config.config import settings


def determine_file_type(filename: str) -> FileType:
    """Determine the file type from the filename"""
    extension = filename.split(".")[-1].lower()
    if extension in ["pdf"]:
        return FileType.PDF
    elif extension in ["docx", "doc"]:
        return FileType.DOCX
    elif extension in ["xlsx", "xls"]:
        return FileType.XLSX
    elif extension in ["pptx", "ppt"]:
        return FileType.PPTX
    elif extension in ["txt"]:
        return FileType.TXT
    else:
        return FileType.UNKNOWN


def save_uploaded_file(file: UploadFile, file_id: str, file_type: FileType) -> str:
    """Save the uploaded file to S3 or local storage"""
    # For development, save to local storage
    if not os.path.exists(settings.UPLOAD_DIR):
        os.makedirs(settings.UPLOAD_DIR)

    file_path = os.path.join(settings.UPLOAD_DIR, f"{file_id}.{file_type.value}")

    with open(file_path, "wb") as buffer:
        content = file.file.read()
        buffer.write(content)
        file.file.seek(0)

    # TODO: Upload to S3 when ready
    # s3_key = f"{file_id}.{file_type.value}"
    # s3_client = boto3.client(
    #     's3',
    #     endpoint_url=settings.S3_ENDPOINT,
    #     aws_access_key_id=settings.S3_ACCESS_KEY,
    #     aws_secret_access_key=settings.S3_SECRET_KEY,
    #     region_name=settings.S3_REGION
    # )
    # s3_client.upload_file(file_path, settings.S3_BUCKET_NAME, s3_key)

    return file_path


def parse_pdf(file_path: str) -> Tuple[List[str], bool, int]:
    """Parse a PDF file and return its content, whether it has images, and page count"""
    doc = fitz.open(file_path)
    page_count = len(doc)
    has_images = False
    content = []

    for page_num, page in enumerate(doc):
        text = page.get_text()
        content.append(text)

        # Check if page has images
        if not has_images and page.get_images():
            has_images = True

    return content, has_images, page_count


def parse_docx(file_path: str) -> Tuple[List[str], bool, int]:
    """Parse a DOCX file and return its content, whether it has images, and page count"""
    doc = docx.Document(file_path)
    content = []
    has_images = False

    # Extract text
    for para in doc.paragraphs:
        content.append(para.text)

    # Check for images
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            has_images = True
            break

    # Approximate page count (not accurate)
    page_count = max(1, len(content) // 40)

    return content, has_images, page_count


def parse_xlsx(file_path: str) -> Tuple[List[str], bool, int]:
    """Parse an XLSX file and return its content, whether it has images, and page count"""
    excel_file = pd.ExcelFile(file_path)
    content = []
    has_images = False  # Excel files typically don't have embedded images we can detect easily

    for sheet_name in excel_file.sheet_names:
        df = pd.read_excel(excel_file, sheet_name=sheet_name)
        content.append(f"Sheet: {sheet_name}")
        content.append(df.to_string(index=False))

    # Each sheet counts as a page
    page_count = len(excel_file.sheet_names)

    return content, has_images, page_count


def parse_pptx(file_path: str) -> Tuple[List[str], bool, int]:
    """Parse a PPTX file and return its content, whether it has images, and page count"""
    prs = Presentation(file_path)
    content = []
    has_images = False

    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
            if not has_images and hasattr(shape, "shape_type") and shape.shape_type == 13:  # Picture
                has_images = True

        content.append(" ".join(slide_text))

    page_count = len(prs.slides)

    return content, has_images, page_count


def parse_txt(file_path: str) -> Tuple[List[str], bool, int]:
    """Parse a TXT file and return its content, whether it has images, and page count"""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()

    # Split by lines and group into pages (approx. 50 lines per page)
    lines = text.split("\n")
    content = []
    for i in range(0, len(lines), 50):
        page_content = "\n".join(lines[i:i+50])
        content.append(page_content)

    if not content:
        content = [text]

    has_images = False  # Text files don't have images
    page_count = max(1, len(content))

    return content, has_images, page_count


def parse_file(file_path: str, file_type: FileType) -> Tuple[List[str], bool, int]:
    """Parse a file based on its type"""
    if file_type == FileType.PDF:
        return parse_pdf(file_path)
    elif file_type == FileType.DOCX:
        return parse_docx(file_path)
    elif file_type == FileType.XLSX:
        return parse_xlsx(file_path)
    elif file_type == FileType.PPTX:
        return parse_pptx(file_path)
    elif file_type == FileType.TXT:
        return parse_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")


def chunk_text(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
    """Split text into chunks with overlap"""
    if not text:
        return []

    chunks = []
    start = 0
    text_length = len(text)

    while start < text_length:
        end = min(start + chunk_size, text_length)

        # Try to find a good breaking point (newline or space)
        if end < text_length:
            # Look for newline first
            newline_pos = text.rfind("\n", start, end)
            if newline_pos > start:
                end = newline_pos + 1
            else:
                # Look for space
                space_pos = text.rfind(" ", start, end)
                if space_pos > start:
                    end = space_pos + 1

        # Add the chunk
        chunks.append(text[start:end])

        # Move the start position, considering overlap
        start = end - chunk_overlap if end < text_length else text_length

    return chunks


def create_chunks_from_content(file_id: str, content: List[str]) -> List[Chunk]:
    """Create chunks from file content"""
    chunks = []
    chunk_index = 0

    for page_num, page_content in enumerate(content):
        page_chunks = chunk_text(
            page_content,
            chunk_size=settings.CHUNK_SIZE,
            chunk_overlap=settings.CHUNK_OVERLAP
        )

        for chunk_text in page_chunks:
            chunk = Chunk(
                id=str(uuid.uuid4()),
                file_id=file_id,
                content=chunk_text,
                page_number=page_num + 1,
                chunk_index=chunk_index,
                created_at=datetime.now(),
                metadata={
                    "page_number": page_num + 1,
                    "chunk_index": chunk_index
                }
            )
            chunks.append(chunk)
            chunk_index += 1

    return chunks
