"""
LlamaIndex service for document processing, indexing, and querying.
"""
import os
import uuid
import time
import asyncio
from typing import List, Dict, Any, Optional
from datetime import datetime
import logging
from enum import Enum

# LlamaIndex imports - using modular package structure
from llama_index.core import (
    VectorStoreIndex,
    SimpleDirectoryReader,
    Document,
    StorageContext,
    Settings,
)
from llama_index.core.retrievers import VectorIndexRetriever
from llama_index.core.node_parser import SentenceSplitter
from llama_index.core.schema import TextNode
from llama_index.core.node_parser import SimpleNodeParser
from llama_index.core.ingestion import IngestionPipeline
from llama_index.embeddings.openai import OpenAIEmbedding
from llama_index.llms.openai import OpenAI
from llama_index.vector_stores.weaviate import WeaviateVectorStore
import weaviate
from weaviate.classes.init import Auth, AdditionalConfig, Timeout

# FastAPI imports
from fastapi import UploadFile, HTTPException

# Local imports
from app.models.db_models import FileType, FileStatus, Chunk
from config.config import settings

# Configure logging
logger = logging.getLogger(__name__)


class ChunkingStrategy(str, Enum):
    """Chunking strategies for document processing."""
    FIXED_SIZE = "fixed_size"
    SEMANTIC = "semantic"
    HYBRID = "hybrid"


class LlamaIndexService:
    """Service for document processing using LlamaIndex."""

    def __init__(self):
        """Initialize the LlamaIndex service."""
        # Configure LlamaIndex settings
        Settings.llm = OpenAI(
            model=settings.DEFAULT_MODEL,
            api_key=settings.OPENAI_API_KEY,
            temperature=0.1
        )
        Settings.embed_model = OpenAIEmbedding(
            model_name=settings.EMBEDDING_MODEL,
            api_key=settings.OPENAI_API_KEY,
            embed_batch_size=10,  # Process 10 chunks at a time to avoid rate limits
        )

        # Initialize Weaviate client if configured
        self.weaviate_client = None
        self.vector_store = None
        if settings.WEAVIATE_URL and settings.WEAVIATE_API_KEY:
            try:
                # Connect to Weaviate cloud using the updated client API
                # Skip initialization checks to avoid gRPC issues
                # Make sure we're using the REST endpoint, not gRPC
                weaviate_url = settings.WEAVIATE_URL
                if not weaviate_url.startswith("https://"):
                    weaviate_url = f"https://{weaviate_url}"

                logger.info(f"Connecting to Weaviate at {weaviate_url}")

                # Use a try-except block with multiple connection attempts
                max_retries = 3
                retry_count = 0
                connection_successful = False

                while retry_count < max_retries and not connection_successful:
                    try:
                        self.weaviate_client = weaviate.connect_to_weaviate_cloud(
                            cluster_url=weaviate_url,
                            auth_credentials=Auth.api_key(settings.WEAVIATE_API_KEY),
                            skip_init_checks=True,  # Skip initialization checks
                            additional_config=AdditionalConfig(
                                timeout=Timeout(
                                    init=settings.WEAVIATE_BATCH_TIMEOUT,  # Increase timeout for initialization
                                    query=settings.WEAVIATE_BATCH_TIMEOUT,  # Increase timeout for queries
                                    batch=settings.WEAVIATE_BATCH_TIMEOUT   # Increase timeout for batch operations
                                )
                            )
                        )
                        connection_successful = True
                        logger.info("Successfully connected to Weaviate")
                    except Exception as e:
                        retry_count += 1
                        logger.warning(f"Weaviate connection attempt {retry_count} failed: {str(e)}")
                        if retry_count < max_retries:
                            logger.info(f"Retrying connection to Weaviate ({retry_count}/{max_retries})...")
                            time.sleep(1)  # Wait 1 second before retrying

                # Only proceed with vector store creation if connection was successful
                if connection_successful and self.weaviate_client:
                    try:
                        # Create vector store with the updated API
                        self.vector_store = WeaviateVectorStore(
                            weaviate_client=self.weaviate_client,
                            index_name=settings.LLAMAINDEX_INDEX_NAME,
                            text_key="text",
                            metadata_keys=["file_id", "user_id", "session_id", "page_number", "chunk_index", "heading", "chunking_strategy"]
                        )

                        # Create schema if it doesn't exist
                        self._create_schema_if_not_exists()
                    except Exception as e:
                        logger.error(f"Error creating vector store: {str(e)}")
                        self.vector_store = None
            except Exception as e:
                logger.error(f"Error connecting to Weaviate: {str(e)}")
                self.weaviate_client = None
                self.vector_store = None

        # Set global settings instead of using ServiceContext (which is deprecated)
        Settings.chunk_size = settings.LLAMAINDEX_CHUNK_SIZE
        Settings.chunk_overlap = settings.LLAMAINDEX_CHUNK_OVERLAP

        # Create storage context if vector store is available
        self.storage_context = None
        if self.vector_store:
            self.storage_context = StorageContext.from_defaults(
                vector_store=self.vector_store
            )

    def get_collection_name_for_user(self, user_id: str) -> str:
        """
        Get the collection name for a specific user.

        Args:
            user_id: The user ID

        Returns:
            The collection name for the user
        """
        if not user_id:
            return settings.LLAMAINDEX_INDEX_NAME

        # Create a user-specific collection name
        # Weaviate doesn't allow underscores in class names, so we'll replace them with hyphens
        # Also, we'll use a shorter version of the user_id to avoid exceeding length limits
        short_user_id = user_id.replace("-", "")[:8]
        return f"{settings.LLAMAINDEX_INDEX_NAME}{short_user_id}"

    def get_vector_store_for_user(self, user_id: str):
        """
        Get a user-specific vector store.

        Args:
            user_id: The user ID

        Returns:
            WeaviateVectorStore instance for the user
        """
        if not self.weaviate_client:
            return None

        # Get user-specific collection name
        collection_name = self.get_collection_name_for_user(user_id)

        # Create vector store with user-specific collection
        return WeaviateVectorStore(
            weaviate_client=self.weaviate_client,
            index_name=collection_name,
            text_key="text",
            metadata_keys=["file_id", "user_id", "session_id", "page_number", "chunk_index", "heading", "chunking_strategy"]
        )

    def _create_user_schema_if_not_exists(self, user_id: str):
        """
        Create Weaviate schema for a specific user if it doesn't exist.

        Args:
            user_id: The user ID
        """
        if not self.weaviate_client:
            return

        collection_name = self.get_collection_name_for_user(user_id)

        try:
            # Check if the collection exists
            try:
                collections = self.weaviate_client.collections.list_all()
                collection_names = []
                for collection in collections:
                    if hasattr(collection, 'name'):
                        collection_names.append(collection.name)
                    elif isinstance(collection, str):
                        collection_names.append(collection)
                    elif isinstance(collection, dict) and 'name' in collection:
                        collection_names.append(collection['name'])
            except Exception as e:
                logger.error(f"Error listing collections: {str(e)}")
                collection_names = []

            # Create collection if it doesn't exist
            if collection_name not in collection_names:
                # Create a new collection
                self.weaviate_client.collections.create(
                    name=collection_name,
                    description="Document chunks for semantic search",
                    vectorizer_config=None,  # We'll provide our own vectors
                    properties=[
                        {
                            "name": "text",
                            "dataType": ["text"],
                            "description": "The text content of the chunk"
                        },
                        {
                            "name": "file_id",
                            "dataType": ["text"],
                            "description": "The ID of the file this chunk belongs to"
                        },
                        {
                            "name": "user_id",
                            "dataType": ["text"],
                            "description": "The ID of the user who owns this chunk"
                        },
                        {
                            "name": "session_id",
                            "dataType": ["text"],
                            "description": "The ID of the session this chunk is associated with"
                        },
                        {
                            "name": "page_number",
                            "dataType": ["int"],
                            "description": "The page number this chunk is from"
                        },
                        {
                            "name": "chunk_index",
                            "dataType": ["int"],
                            "description": "The index of this chunk within the file"
                        },
                        {
                            "name": "heading",
                            "dataType": ["text"],
                            "description": "The heading or title of the section"
                        },
                        {
                            "name": "chunking_strategy",
                            "dataType": ["text"],
                            "description": "The chunking strategy used (fixed_size, semantic, hybrid)"
                        },
                        {
                            "name": "metadata",
                            "dataType": ["text"],
                            "description": "Additional metadata about the chunk"
                        }
                    ]
                )
                logger.info(f"Created collection {collection_name} in Weaviate")
        except Exception as e:
            logger.error(f"Error creating Weaviate schema for user {user_id}: {str(e)}")

    def _create_schema_if_not_exists(self):
        """Create Weaviate schema if it doesn't exist."""
        if not self.weaviate_client:
            return

        try:
            # Check if the collection exists
            try:
                collections = self.weaviate_client.collections.list_all()
                collection_names = []
                for collection in collections:
                    if hasattr(collection, 'name'):
                        collection_names.append(collection.name)
                    elif isinstance(collection, str):
                        collection_names.append(collection)
                    elif isinstance(collection, dict) and 'name' in collection:
                        collection_names.append(collection['name'])
            except Exception as e:
                logger.error(f"Error listing collections: {str(e)}")
                collection_names = []

            # Create collection if it doesn't exist
            if settings.LLAMAINDEX_INDEX_NAME not in collection_names:
                # Create a new collection
                self.weaviate_client.collections.create(
                    name=settings.LLAMAINDEX_INDEX_NAME,
                    description="Document chunks for semantic search",
                    vectorizer_config=None,  # We'll provide our own vectors
                    properties=[
                        {
                            "name": "text",
                            "dataType": ["text"],
                            "description": "The text content of the chunk"
                        },
                        {
                            "name": "file_id",
                            "dataType": ["text"],
                            "description": "The ID of the file this chunk belongs to"
                        },
                        {
                            "name": "user_id",
                            "dataType": ["text"],
                            "description": "The ID of the user who owns this chunk"
                        },
                        {
                            "name": "session_id",
                            "dataType": ["text"],
                            "description": "The ID of the session this chunk is associated with"
                        },
                        {
                            "name": "page_number",
                            "dataType": ["int"],
                            "description": "The page number this chunk is from"
                        },
                        {
                            "name": "chunk_index",
                            "dataType": ["int"],
                            "description": "The index of this chunk within the file"
                        },
                        {
                            "name": "heading",
                            "dataType": ["text"],
                            "description": "The heading or title of the section"
                        },
                        {
                            "name": "chunking_strategy",
                            "dataType": ["text"],
                            "description": "The chunking strategy used (fixed_size, semantic, hybrid)"
                        },
                        {
                            "name": "metadata",
                            "dataType": ["text"],
                            "description": "Additional metadata about the chunk"
                        }
                    ]
                )
                logger.info(f"Created collection {settings.LLAMAINDEX_INDEX_NAME} in Weaviate")
        except Exception as e:
            logger.error(f"Error creating Weaviate schema: {str(e)}")

    async def process_file(self, file_path: str, file_id: str, user_id: str,
                          file_type: FileType, chunking_strategy: ChunkingStrategy = ChunkingStrategy.HYBRID,
                          session_id: Optional[str] = None) -> Dict[str, Any]:
        """
        Process a file using LlamaIndex.

        Args:
            file_path: Path to the file
            file_id: Unique ID for the file
            user_id: ID of the user who uploaded the file
            file_type: Type of the file
            chunking_strategy: Chunking strategy to use
            session_id: ID of the session (optional)

        Returns:
            Dict containing processing results
        """
        try:
            # Validate file path
            if not self._validate_file_path(file_path):
                raise FileNotFoundError(f"File not found or not readable: {file_path}")

            # If file_type is UNKNOWN, try to determine it from the file path
            if file_type == FileType.UNKNOWN:
                detected_type = self._determine_file_type(file_path)
                if detected_type != FileType.UNKNOWN:
                    file_type = detected_type
                    logger.info(f"Detected file type {file_type} for {file_path}")

            # Load the document
            documents = await self._load_document(file_path, file_type)

            if not documents:
                raise ValueError(f"No content could be extracted from {file_path}")

            # Check if document has images
            has_images = self._check_for_images(documents, file_type)

            # Get page count
            page_count = len(documents)

            # Create nodes with appropriate chunking strategy
            nodes = await self._create_nodes(documents, file_id, user_id, chunking_strategy, session_id)

            if not nodes:
                raise ValueError(f"No chunks could be created from {file_path}")

            # Create user-specific vector store and storage context
            if self.weaviate_client:
                # Create schema for user if it doesn't exist
                self._create_user_schema_if_not_exists(user_id)

                # Get user-specific vector store
                user_vector_store = self.get_vector_store_for_user(user_id)

                if user_vector_store:
                    # Use user-specific vector store with batched processing
                    await self._store_nodes_in_batches_for_user(nodes, user_vector_store)
                else:
                    # Use in-memory index if no vector store
                    VectorStoreIndex(
                        nodes=nodes,
                    )
            else:
                # Use in-memory index if no Weaviate client
                VectorStoreIndex(
                    nodes=nodes,
                )

            # Create chunks for database storage
            chunks = self._create_chunks_from_nodes(nodes, file_id)

            return {
                "file_id": file_id,
                "status": "processed",
                "page_count": page_count,
                "has_images": has_images,
                "chunk_count": len(chunks),
                "chunks": chunks
            }

        except Exception as e:
            logger.error(f"Error processing file {file_id}: {str(e)}")
            raise

    async def _load_document(self, file_path: str, file_type: FileType) -> List[Document]:
        """
        Load a document using proper file type parsing.

        Args:
            file_path: Path to the file
            file_type: Type of the file

        Returns:
            List of Document objects
        """
        try:
            # Validate file exists and is readable
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"File not found: {file_path}")

            # Log file information for debugging
            file_size = os.path.getsize(file_path)
            logger.info(f"Processing file: {file_path}, size: {file_size} bytes, type: {file_type}")

            # Check if file is actually readable
            try:
                with open(file_path, "rb") as f:
                    first_bytes = f.read(10)
                    logger.info(f"First 10 bytes of file: {first_bytes}")
            except Exception as read_error:
                logger.error(f"Cannot read file {file_path}: {str(read_error)}")
                raise

            # Use proper file type parsing instead of SimpleDirectoryReader
            # to ensure we get readable text content, not raw file structure
            text_content = ""

            if file_type == FileType.PDF:
                # Use PyPDF2 for proper PDF text extraction
                from pypdf import PdfReader
                try:
                    logger.info(f"Attempting to read PDF file: {file_path}")
                    with open(file_path, "rb") as f:
                        # Check if file starts with PDF header
                        f.seek(0)
                        header = f.read(4)
                        if header != b'%PDF':
                            logger.error(f"File {file_path} does not appear to be a valid PDF (header: {header})")
                            raise ValueError(f"File does not appear to be a valid PDF file")

                        f.seek(0)  # Reset to beginning
                        pdf = PdfReader(f)
                        logger.info(f"PDF has {len(pdf.pages)} pages")

                        for page_num, page in enumerate(pdf.pages):
                            try:
                                page_text = page.extract_text()
                                if page_text.strip():  # Only add non-empty pages
                                    logger.debug(f"Extracted {len(page_text)} characters from page {page_num + 1}")
                                    text_content += f"Page {page_num + 1}:\n{page_text}\n\n"
                                else:
                                    logger.warning(f"Page {page_num + 1} appears to be empty or contains no extractable text")
                            except Exception as page_error:
                                logger.warning(f"Error extracting text from page {page_num + 1}: {str(page_error)}")
                                continue

                        logger.info(f"Total extracted text length: {len(text_content)} characters")
                except Exception as pdf_error:
                    logger.error(f"Error reading PDF file {file_path}: {str(pdf_error)}")
                    # Try to provide more specific error information
                    try:
                        with open(file_path, "rb") as f:
                            first_100_bytes = f.read(100)
                            logger.error(f"First 100 bytes of file: {first_100_bytes}")
                    except:
                        pass
                    raise ValueError(f"Unable to read PDF file: {str(pdf_error)}")

            elif file_type == FileType.DOCX:
                # Use python-docx for proper DOCX text extraction
                from docx import Document as DocxDocument
                try:
                    doc = DocxDocument(file_path)
                    for para in doc.paragraphs:
                        if para.text.strip():  # Only add non-empty paragraphs
                            text_content += para.text + "\n"
                except Exception as docx_error:
                    logger.error(f"Error reading DOCX file {file_path}: {str(docx_error)}")
                    raise ValueError(f"Unable to read DOCX file: {str(docx_error)}")

            elif file_type == FileType.XLSX:
                # Use openpyxl for proper Excel text extraction
                from openpyxl import load_workbook
                try:
                    wb = load_workbook(file_path)
                    for sheet_name in wb.sheetnames:
                        ws = wb[sheet_name]
                        text_content += f"Sheet: {sheet_name}\n"
                        for row in ws.iter_rows(values_only=True):
                            row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                            if row_text.strip():  # Only add non-empty rows
                                text_content += row_text + "\n"
                        text_content += "\n"
                except Exception as xlsx_error:
                    logger.error(f"Error reading XLSX file {file_path}: {str(xlsx_error)}")
                    raise ValueError(f"Unable to read XLSX file: {str(xlsx_error)}")

            elif file_type == FileType.PPTX:
                # Use python-pptx for proper PowerPoint text extraction
                from pptx import Presentation
                try:
                    prs = Presentation(file_path)
                    for slide_num, slide in enumerate(prs.slides):
                        text_content += f"Slide {slide_num + 1}:\n"
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                text_content += shape.text + "\n"
                        text_content += "\n"
                except Exception as pptx_error:
                    logger.error(f"Error reading PPTX file {file_path}: {str(pptx_error)}")
                    raise ValueError(f"Unable to read PPTX file: {str(pptx_error)}")

            elif file_type == FileType.TXT:
                # Handle text files with proper encoding detection
                try:
                    # Try UTF-8 first
                    with open(file_path, "r", encoding="utf-8") as f:
                        text_content = f.read()
                except UnicodeDecodeError:
                    # Fallback to other encodings
                    try:
                        with open(file_path, "r", encoding="latin-1") as f:
                            text_content = f.read()
                    except Exception as txt_error:
                        logger.error(f"Error reading TXT file {file_path}: {str(txt_error)}")
                        raise ValueError(f"Unable to read TXT file: {str(txt_error)}")

            else:
                # Unknown file type - try to read as text with error handling
                logger.warning(f"Unknown file type {file_type} for {file_path}, attempting text extraction")
                try:
                    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                        text_content = f.read()
                except Exception as unknown_error:
                    logger.error(f"Error reading unknown file type {file_path}: {str(unknown_error)}")
                    raise ValueError(f"Unable to read file of type {file_type}: {str(unknown_error)}")

            # Create a single document with all the extracted text
            if not text_content.strip():
                logger.warning(f"No text content extracted from {file_path}")
                text_content = "No readable text content found in this document."
            else:
                # Validate that the text content is actually readable (not binary garbage)
                try:
                    # Check if text contains mostly printable characters
                    printable_chars = sum(1 for c in text_content[:1000] if c.isprintable() or c.isspace())
                    total_chars = len(text_content[:1000])
                    if total_chars > 0:
                        printable_ratio = printable_chars / total_chars
                        logger.info(f"Text content printable ratio: {printable_ratio:.2f}")

                        if printable_ratio < 0.7:  # Less than 70% printable characters
                            logger.error(f"Text content appears to be corrupted (only {printable_ratio:.2f} printable)")
                            logger.error(f"Sample of extracted text: {repr(text_content[:200])}")
                            raise ValueError("Extracted text appears to be corrupted or binary data")

                    logger.info(f"Successfully extracted {len(text_content)} characters of readable text")
                except Exception as validation_error:
                    logger.error(f"Text validation failed: {str(validation_error)}")
                    raise

            documents = [Document(
                text=text_content,
                metadata={
                    "file_path": file_path,
                    "file_type": file_type.value,
                    "file_name": os.path.basename(file_path)
                }
            )]

            # For PDF files, try to split by page markers if they exist
            if file_type == FileType.PDF and len(documents) == 1:
                text = documents[0].text
                # Look for page markers that we added during extraction
                if "Page " in text:
                    # Split by page markers
                    pages = text.split("Page ")[1:]  # Skip the first empty element
                    if len(pages) > 1:
                        # Create a document for each page
                        documents = []
                        for i, page_content in enumerate(pages):
                            # Add back the page marker
                            page_text = f"Page {page_content}"
                            documents.append(Document(
                                text=page_text,
                                metadata={
                                    "page_number": i + 1,
                                    "file_path": file_path,
                                    "file_type": file_type.value,
                                    "file_name": os.path.basename(file_path)
                                }
                            ))
                    else:
                        # Add page metadata to the single document
                        documents[0].metadata["page_number"] = 1
                else:
                    # Add page metadata to the single document
                    documents[0].metadata["page_number"] = 1
            else:
                # Add page metadata to the single document
                if documents:
                    documents[0].metadata["page_number"] = 1

            return documents

        except Exception as e:
            logger.error(f"Error loading document {file_path}: {str(e)}")
            raise

    def _check_for_images(self, _docs: List[Document], file_type: FileType) -> bool:
        """
        Check if a document contains images.

        Args:
            _docs: List of Document objects (unused)
            file_type: Type of the file

        Returns:
            True if the document contains images, False otherwise
        """
        # For now, we'll use a simple heuristic based on file type
        # In a more advanced implementation, we could analyze the document content
        if file_type in [FileType.PDF, FileType.DOCX, FileType.PPTX]:
            # These file types commonly contain images
            # For a more accurate check, we would need to parse the document structure
            return True

        return False

    async def _create_nodes(self, documents: List[Document], file_id: str, user_id: str,
                           chunking_strategy: ChunkingStrategy, session_id: Optional[str] = None) -> List[TextNode]:
        """
        Create nodes from documents using the specified chunking strategy.

        Args:
            documents: List of Document objects
            file_id: Unique ID for the file
            user_id: ID of the user who uploaded the file
            chunking_strategy: Chunking strategy to use
            session_id: ID of the session (optional)

        Returns:
            List of TextNode objects
        """
        nodes = []

        if chunking_strategy == ChunkingStrategy.FIXED_SIZE:
            # Use simple fixed-size chunking
            node_parser = SimpleNodeParser.from_defaults(
                chunk_size=settings.LLAMAINDEX_CHUNK_SIZE,
                chunk_overlap=settings.LLAMAINDEX_CHUNK_OVERLAP
            )
        elif chunking_strategy == ChunkingStrategy.SEMANTIC:
            # Use sentence-based chunking for more semantic coherence
            node_parser = SentenceSplitter(
                chunk_size=settings.LLAMAINDEX_CHUNK_SIZE,
                chunk_overlap=settings.LLAMAINDEX_CHUNK_OVERLAP,
                paragraph_separator="\n\n",
                secondary_chunking_regex="[^,.;。]+[,.;。]?",
            )
        else:  # ChunkingStrategy.HYBRID (default)
            # Use a combination based on document type
            # For now, we'll use the same as SEMANTIC, but this could be enhanced
            node_parser = SentenceSplitter(
                chunk_size=settings.LLAMAINDEX_CHUNK_SIZE,
                chunk_overlap=settings.LLAMAINDEX_CHUNK_OVERLAP,
                paragraph_separator="\n\n",
                secondary_chunking_regex="[^,.;。]+[,.;。]?",
            )

        # Create an ingestion pipeline
        pipeline = IngestionPipeline(
            transformations=[node_parser],
        )

        # Process documents through the pipeline
        for doc_idx, doc in enumerate(documents):
            # Add file and user metadata to the document
            doc.metadata["file_id"] = file_id
            doc.metadata["user_id"] = user_id
            doc.metadata["chunking_strategy"] = chunking_strategy

            # Add session_id if provided
            if session_id:
                doc.metadata["session_id"] = session_id

            # If page_number is not set, use the document index
            if "page_number" not in doc.metadata:
                doc.metadata["page_number"] = doc_idx + 1

        # Run the pipeline
        nodes = pipeline.run(documents)

        # Add additional metadata to nodes
        for i, node in enumerate(nodes):
            node.metadata["chunk_index"] = i

            # Try to extract heading from the text
            lines = node.text.split("\n")
            if lines and len(lines[0]) < 100:  # Simple heuristic for headings
                node.metadata["heading"] = lines[0]
            else:
                node.metadata["heading"] = None

        return nodes

    async def _store_nodes_in_batches_for_user(self, nodes: List[TextNode], user_vector_store) -> None:
        """
        Store nodes in Weaviate using batched processing for a specific user.

        Args:
            nodes: List of TextNode objects to store
            user_vector_store: User-specific vector store
        """
        if not user_vector_store or not self.weaviate_client:
            logger.warning("Vector store not available, skipping batch storage")
            return

        try:
            # Get the total number of nodes
            total_nodes = len(nodes)
            logger.info(f"Starting batch processing of {total_nodes} nodes for user")

            # Calculate number of batches
            batch_size = settings.WEAVIATE_BATCH_SIZE
            num_batches = (total_nodes + batch_size - 1) // batch_size  # Ceiling division

            # Process nodes in batches
            for batch_idx in range(num_batches):
                start_idx = batch_idx * batch_size
                end_idx = min(start_idx + batch_size, total_nodes)
                batch_nodes = nodes[start_idx:end_idx]

                logger.info(f"Processing batch {batch_idx + 1}/{num_batches} with {len(batch_nodes)} nodes")

                # Create a temporary storage context for this batch
                batch_storage_context = StorageContext.from_defaults(
                    vector_store=user_vector_store
                )

                # Process the batch with retries
                retry_count = 0
                max_retries = settings.WEAVIATE_MAX_RETRIES
                success = False

                while retry_count < max_retries and not success:
                    try:
                        # Create a temporary index for this batch
                        VectorStoreIndex(
                            nodes=batch_nodes,
                            storage_context=batch_storage_context,
                        )
                        success = True
                        logger.info(f"Successfully processed batch {batch_idx + 1}/{num_batches}")
                    except Exception as e:
                        retry_count += 1
                        logger.warning(f"Batch {batch_idx + 1} attempt {retry_count} failed: {str(e)}")
                        if retry_count < max_retries:
                            logger.info(f"Retrying batch {batch_idx + 1} (attempt {retry_count + 1}/{max_retries})...")
                            # Wait before retrying with exponential backoff
                            await asyncio.sleep(2 ** retry_count)
                        else:
                            logger.error(f"Failed to process batch {batch_idx + 1} after {max_retries} attempts")
                            # Continue with next batch instead of failing the entire process

            logger.info(f"Completed batch processing of {total_nodes} nodes for user")
        except Exception as e:
            logger.error(f"Error in batch processing for user: {str(e)}")
            # Don't raise the exception to allow the process to continue
            # The document will be marked as processed even if some batches failed

    async def _store_nodes_in_batches(self, nodes: List[TextNode]) -> None:
        """
        Store nodes in Weaviate using batched processing to avoid timeouts.

        Args:
            nodes: List of TextNode objects to store
        """
        if not self.vector_store or not self.weaviate_client:
            logger.warning("Vector store not available, skipping batch storage")
            return

        try:
            # Get the total number of nodes
            total_nodes = len(nodes)
            logger.info(f"Starting batch processing of {total_nodes} nodes")

            # Calculate number of batches
            batch_size = settings.WEAVIATE_BATCH_SIZE
            num_batches = (total_nodes + batch_size - 1) // batch_size  # Ceiling division

            # Process nodes in batches
            for batch_idx in range(num_batches):
                start_idx = batch_idx * batch_size
                end_idx = min(start_idx + batch_size, total_nodes)
                batch_nodes = nodes[start_idx:end_idx]

                logger.info(f"Processing batch {batch_idx + 1}/{num_batches} with {len(batch_nodes)} nodes")

                # Create a temporary storage context for this batch
                batch_storage_context = StorageContext.from_defaults(
                    vector_store=self.vector_store
                )

                # Process the batch with retries
                retry_count = 0
                max_retries = settings.WEAVIATE_MAX_RETRIES
                success = False

                while retry_count < max_retries and not success:
                    try:
                        # Create a temporary index for this batch
                        VectorStoreIndex(
                            nodes=batch_nodes,
                            storage_context=batch_storage_context,
                        )
                        success = True
                        logger.info(f"Successfully processed batch {batch_idx + 1}/{num_batches}")
                    except Exception as e:
                        retry_count += 1
                        logger.warning(f"Batch {batch_idx + 1} attempt {retry_count} failed: {str(e)}")
                        if retry_count < max_retries:
                            logger.info(f"Retrying batch {batch_idx + 1} (attempt {retry_count + 1}/{max_retries})...")
                            # Wait before retrying with exponential backoff
                            await asyncio.sleep(2 ** retry_count)
                        else:
                            logger.error(f"Failed to process batch {batch_idx + 1} after {max_retries} attempts")
                            # Continue with next batch instead of failing the entire process

            logger.info(f"Completed batch processing of {total_nodes} nodes")
        except Exception as e:
            logger.error(f"Error in batch processing: {str(e)}")
            # Don't raise the exception to allow the process to continue
            # The document will be marked as processed even if some batches failed

    def _create_chunks_from_nodes(self, nodes: List[TextNode], file_id: str) -> List[Chunk]:
        """
        Create Chunk objects from TextNode objects for database storage.

        Args:
            nodes: List of TextNode objects
            file_id: Unique ID for the file

        Returns:
            List of Chunk objects
        """
        chunks = []

        for node in nodes:
            # Create a Chunk object from the node
            chunk = Chunk(
                id=str(uuid.uuid4()),
                file_id=file_id,
                content=node.text,
                page_number=node.metadata.get("page_number"),
                chunk_index=node.metadata.get("chunk_index", 0),
                embedding_id=node.id_,  # Use the node ID as the embedding ID
                created_at=datetime.now(),
                metadata={
                    "page_number": node.metadata.get("page_number"),
                    "chunk_index": node.metadata.get("chunk_index", 0),
                    "heading": node.metadata.get("heading"),
                    "chunking_strategy": node.metadata.get("chunking_strategy"),
                    "file_id": file_id,
                    "user_id": node.metadata.get("user_id"),
                }
            )
            chunks.append(chunk)

        return chunks

    async def get_document_chunks(self, file_id: str, user_id: str, limit: int = 3) -> Dict[str, Any]:
        """
        Get representative chunks from a document.

        Args:
            file_id: ID of the file to get chunks from
            user_id: ID of the user who owns the file
            limit: Maximum number of chunks to return

        Returns:
            Dict containing document chunks
        """
        try:
            if not self.weaviate_client:
                raise HTTPException(status_code=500, detail="Vector store not configured")

            # Get user-specific collection name
            collection_name = self.get_collection_name_for_user(user_id)

            # Create a filter for the specified file ID and user ID
            filter_dict = {
                "operator": "And",
                "operands": [
                    {
                        "path": ["file_id"],
                        "operator": "Equal",
                        "valueString": file_id
                    },
                    {
                        "path": ["user_id"],
                        "operator": "Equal",
                        "valueString": user_id
                    }
                ]
            }

            # Get chunks from the vector store
            try:
                # Try using the v4 API first
                collection = self.weaviate_client.collections.get(collection_name)

                # Query for chunks
                results = collection.query.fetch_objects(
                    limit=limit,
                    filters=filter_dict,
                    include_vector=False
                )

                # Extract chunks
                chunks = []
                for obj in results.objects:
                    chunks.append({
                        "content": obj.properties.get("content", ""),
                        "file_id": obj.properties.get("file_id", ""),
                        "page_number": obj.properties.get("page_number", 0),
                        "chunk_index": obj.properties.get("chunk_index", 0),
                        "heading": obj.properties.get("heading", ""),
                        "metadata": obj.properties.get("metadata", {})
                    })
            except (AttributeError, Exception) as e:
                logger.error(f"Error using v4 API to get chunks: {str(e)}")
                # Fall back to a simpler approach - create index and retriever with user-specific vector store
                user_vector_store = self.get_vector_store_for_user(user_id)
                if user_vector_store:
                    storage_context = StorageContext.from_defaults(vector_store=user_vector_store)
                    index = VectorStoreIndex.from_vector_store(
                        vector_store=user_vector_store,
                        storage_context=storage_context
                    )
                    retriever = VectorIndexRetriever(
                        index=index,
                        similarity_top_k=limit
                    )

                    # Use a generic query to get some representative chunks
                    nodes = retriever.retrieve("summary of this document")

                    # Extract chunks
                    chunks = []
                    for node in nodes:
                        chunks.append({
                            "content": node.text,
                            "file_id": node.metadata.get("file_id", ""),
                            "page_number": node.metadata.get("page_number", 0),
                            "chunk_index": node.metadata.get("chunk_index", 0),
                            "heading": node.metadata.get("heading", ""),
                            "metadata": node.metadata
                        })
                else:
                    # No vector store available
                    chunks = []

            return {
                "file_id": file_id,
                "user_id": user_id,
                "chunks": chunks
            }

        except Exception as e:
            logger.error(f"Error getting document chunks: {str(e)}")

            # Provide helpful error message
            error_message = self._get_helpful_error_message(str(e))

            return {
                "file_id": file_id,
                "user_id": user_id,
                "chunks": [],
                "error": error_message
            }

    async def query_documents(self, query: str, file_ids: List[str], user_id: str,
                             top_k: int = 5, session_id: Optional[str] = None) -> Dict[str, Any]:
        """
        Query documents using LlamaIndex.

        Args:
            query: Query string
            file_ids: List of file IDs to search in
            user_id: ID of the user making the query
            top_k: Number of results to return
            session_id: ID of the session (optional, for prioritizing session-specific chunks)

        Returns:
            Dict containing query results
        """
        try:
            if not self.weaviate_client:
                raise HTTPException(status_code=500, detail="Vector store not configured")

            # Get user-specific vector store
            user_vector_store = self.get_vector_store_for_user(user_id)
            if not user_vector_store:
                raise HTTPException(status_code=500, detail="User vector store not available")

            # Create storage context and index with user-specific vector store
            storage_context = StorageContext.from_defaults(vector_store=user_vector_store)
            index = VectorStoreIndex.from_vector_store(
                vector_store=user_vector_store,
                storage_context=storage_context
            )

            # Create retriever with simplified approach (no complex filters for now)
            # LlamaIndex filters have compatibility issues, so we'll use basic retrieval
            retriever = VectorIndexRetriever(
                index=index,
                similarity_top_k=top_k
            )

            # Get the query embedding (not needed when using retriever directly)
            # query_embedding = Settings.embed_model.get_query_embedding(query)

            # Execute the query
            nodes = retriever.retrieve(query)

            # Create and return the response
            return self._create_response(query, nodes)
        except Exception as e:
            logger.error(f"Error querying documents: {str(e)}")

            # Return a helpful error response instead of raising
            return {
                "response": self._get_helpful_error_message(str(e)),
                "source_documents": [],
                "model_used": getattr(Settings.llm, 'model', getattr(Settings.llm, 'model_name', 'unknown'))
            }

    def _create_response(self, query: str, nodes: List) -> Dict[str, Any]:
        """
        Create a response from retrieved nodes.

        Args:
            query: The user's query
            nodes: The retrieved nodes

        Returns:
            Dict containing the response
        """
        # Create a response
        llm = Settings.llm
        response_text = llm.complete(
            f"""You are AnyDocAI, an AI document assistant that helps users understand their documents.

            Use the following context from the user's documents to answer their question. If you don't know the answer, say you don't know.
            Don't try to make up an answer. Always be helpful, concise, and professional.

            Context:
            {' '.join([node.text for node in nodes])}

            Question: {query}

            Answer:"""
        ).text

        # Format the response
        response = {
            "response": response_text,
            "source_documents": [
                {
                    "content": node.text,
                    "metadata": node.metadata,
                    "score": node.score if hasattr(node, "score") else None
                }
                for node in nodes
            ],
            "model_used": getattr(Settings.llm, 'model', getattr(Settings.llm, 'model_name', 'unknown'))
        }

        return response

    def _get_helpful_error_message(self, error_str: str) -> str:
        """
        Generate a helpful error message based on the error type.

        Args:
            error_str: The original error string

        Returns:
            A helpful error message for the user
        """
        error_lower = error_str.lower()

        if "connection" in error_lower and ("reset" in error_lower or "unavailable" in error_lower):
            return "Document database connection issue - please try again in a moment"
        elif "weaviate" in error_lower or "vector" in error_lower:
            return "Vector database temporarily unavailable - please try again"
        elif "timeout" in error_lower:
            return "Request timed out - please try again or rephrase your question"
        elif "no documents" in error_lower or "no chunks" in error_lower:
            return "No relevant documents found for this query"
        else:
            return f"Document processing error: {error_str}"

    async def process_uploaded_file(self, file: UploadFile, user_id: str,
                                   chunking_strategy: ChunkingStrategy = ChunkingStrategy.HYBRID) -> Dict[str, Any]:
        """
        Process an uploaded file.

        Args:
            file: Uploaded file
            user_id: ID of the user who uploaded the file
            chunking_strategy: Chunking strategy to use

        Returns:
            Dict containing processing results
        """
        try:
            # Check file size
            file.file.seek(0, os.SEEK_END)
            file_size = file.file.tell()
            file.file.seek(0)

            if file_size > settings.MAX_UPLOAD_SIZE:
                raise HTTPException(status_code=400, detail="File too large")

            # Determine file type
            file_type = self._determine_file_type(file.filename)
            if file_type == FileType.UNKNOWN:
                raise HTTPException(status_code=400, detail="Unsupported file type")

            # Generate a unique ID for the file
            file_id = str(uuid.uuid4())

            # Save the file temporarily
            temp_file_path = os.path.join(settings.UPLOAD_DIR, f"{file_id}.{file_type.value}")
            os.makedirs(os.path.dirname(temp_file_path), exist_ok=True)

            with open(temp_file_path, "wb") as buffer:
                content = file.file.read()
                buffer.write(content)
                file.file.seek(0)

            # Process the file
            result = await self.process_file(
                file_path=temp_file_path,
                file_id=file_id,
                user_id=user_id,
                file_type=file_type,
                chunking_strategy=chunking_strategy
            )

            # Create file record (for future database integration)
            # {
            #     "id": file_id,
            #     "user_id": user_id,
            #     "filename": file.filename,
            #     "original_filename": file.filename,
            #     "file_type": file_type,
            #     "file_size": file_size,
            #     "status": FileStatus.PROCESSED,
            #     "s3_key": temp_file_path,  # For now, just store the local path
            #     "created_at": datetime.now(),
            #     "updated_at": datetime.now(),
            #     "processed_at": datetime.now(),
            #     "page_count": result.get("page_count", 0),
            #     "has_images": result.get("has_images", False)
            # }

            # TODO: Save file record to database

            # Return the result
            return {
                "id": file_id,
                "filename": file.filename,
                "file_type": file_type,
                "file_size": file_size,
                "status": FileStatus.PROCESSED,
                "created_at": datetime.now(),
                "chunks": result.get("chunks", [])
            }

        except Exception as e:
            logger.error(f"Error processing uploaded file: {str(e)}")
            raise

    def _determine_file_type(self, filename: str) -> FileType:
        """
        Determine the file type from the filename.

        Args:
            filename: Name of the file

        Returns:
            FileType enum value
        """
        if not filename:
            return FileType.UNKNOWN

        # Handle both full filenames and just extensions
        if filename.startswith('.'):
            extension = filename[1:].lower()
        else:
            extension = filename.split(".")[-1].lower() if "." in filename else filename.lower()

        if extension in ["pdf"]:
            return FileType.PDF
        elif extension in ["docx", "doc"]:
            return FileType.DOCX
        elif extension in ["xlsx", "xls"]:
            return FileType.XLSX
        elif extension in ["pptx", "ppt"]:
            return FileType.PPTX
        elif extension in ["txt"]:
            return FileType.TXT
        else:
            return FileType.UNKNOWN

    def _validate_file_path(self, file_path: str) -> bool:
        """
        Validate that a file path exists and is readable.

        Args:
            file_path: Path to the file

        Returns:
            True if file is valid, False otherwise
        """
        try:
            return os.path.exists(file_path) and os.path.isfile(file_path) and os.access(file_path, os.R_OK)
        except Exception:
            return False

    def close_connections(self):
        """Close all connections and clean up resources."""
        try:
            if self.weaviate_client:
                logger.info("Closing Weaviate client connection...")
                self.weaviate_client.close()
                self.weaviate_client = None
                logger.info("Weaviate client connection closed successfully")
        except Exception as e:
            logger.error(f"Error closing Weaviate client: {str(e)}")

        # Clear vector store reference
        self.vector_store = None

# Create a singleton instance
llama_index_service = LlamaIndexService()
