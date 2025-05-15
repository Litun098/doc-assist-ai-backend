"""
Document processing service using LlamaIndex.
Handles different file types and implements hybrid chunking strategies.
"""
import os
import time
import logging
import asyncio
from typing import List, Dict, Any, Optional, Tuple
from enum import Enum

# LlamaIndex imports
from llama_index.core import Document as LlamaDocument
from llama_index.core import VectorStoreIndex, StorageContext
from llama_index.core.node_parser import SentenceSplitter
from llama_index.core.schema import TextNode

# Import document processing libraries
from pypdf import PdfReader
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

# Weaviate imports
from llama_index.vector_stores.weaviate import WeaviateVectorStore
import weaviate

# Local imports
from app.models.db_models import FileType
from config.config import settings

# Configure logging
logger = logging.getLogger(__name__)

class ChunkingStrategy(str, Enum):
    """Enum for chunking strategies."""
    FIXED_SIZE = "fixed_size"
    TOPIC_BASED = "topic_based"
    HYBRID = "hybrid"

class DocumentProcessor:
    """Document processor service using LlamaIndex."""

    def __init__(self):
        """Initialize the document processor."""
        # Initialize Weaviate client if configured
        self.weaviate_client = None
        self.use_weaviate = False

        # Initialize Weaviate if configured
        if settings.WEAVIATE_URL and settings.WEAVIATE_API_KEY:
            try:
                from weaviate.classes.init import Auth, AdditionalConfig, Timeout

                # Make sure we're using the REST endpoint, not gRPC
                weaviate_url = settings.WEAVIATE_URL
                if not weaviate_url.startswith("https://"):
                    weaviate_url = f"https://{weaviate_url}"

                logger.info(f"Connecting to Weaviate at {weaviate_url}")
                self.weaviate_client = weaviate.connect_to_weaviate_cloud(
                    cluster_url=weaviate_url,
                    auth_credentials=Auth.api_key(settings.WEAVIATE_API_KEY),
                    skip_init_checks=True,  # Skip initialization checks
                    additional_config=AdditionalConfig(
                        timeout=Timeout(
                            init=settings.WEAVIATE_BATCH_TIMEOUT,  # Increase timeout for initialization
                            query=settings.WEAVIATE_BATCH_TIMEOUT,  # Increase timeout for queries
                            batch=settings.WEAVIATE_BATCH_TIMEOUT   # Increase timeout for batch operations
                        )
                    )
                )
                self.use_weaviate = True
                logger.info("Using Weaviate for vector storage")
            except Exception as e:
                logger.error(f"Error connecting to Weaviate: {str(e)}")
                logger.info("Falling back to in-memory storage")

        # Initialize node parsers for different chunking strategies
        self.fixed_size_parser = SentenceSplitter(
            chunk_size=settings.CHUNK_SIZE,
            chunk_overlap=settings.CHUNK_OVERLAP
        )

        self.topic_based_parser = SentenceSplitter(
            chunk_size=settings.CHUNK_SIZE,
            chunk_overlap=settings.CHUNK_OVERLAP,
            paragraph_separator="\n\n",
            secondary_chunking_regex="|".join(settings.HEADING_PATTERNS)
        )

    def detect_file_type(self, file_path: str) -> FileType:
        """
        Detect the file type based on the file extension.

        Args:
            file_path: Path to the file

        Returns:
            FileType enum value
        """
        _, ext = os.path.splitext(file_path)
        ext = ext.lower().lstrip('.')

        if ext == 'pdf':
            return FileType.PDF
        elif ext in ['docx', 'doc']:
            return FileType.DOCX
        elif ext in ['xlsx', 'xls']:
            return FileType.XLSX
        elif ext in ['pptx', 'ppt']:
            return FileType.PPTX
        elif ext == 'txt':
            return FileType.TXT
        else:
            return FileType.UNKNOWN

    def get_chunking_strategy(self, file_type: FileType) -> ChunkingStrategy:
        """
        Determine the chunking strategy based on file type.

        Args:
            file_type: FileType enum value

        Returns:
            ChunkingStrategy enum value
        """
        if file_type in settings.TOPIC_BASED_FILETYPES:
            return ChunkingStrategy.TOPIC_BASED
        elif file_type in settings.FIXED_SIZE_FILETYPES:
            return ChunkingStrategy.FIXED_SIZE
        else:
            return ChunkingStrategy.HYBRID

    # Document loader methods removed in favor of using SimpleDirectoryReader directly

    def load_document(self, file_path: str, file_type: Optional[FileType] = None) -> List[LlamaDocument]:
        """
        Load a document using the appropriate loader.

        Args:
            file_path: Path to the file
            file_type: FileType enum value (optional, will be detected if not provided)

        Returns:
            List of LlamaIndex Document objects
        """
        try:
            if file_type is None:
                file_type = self.detect_file_type(file_path)

            # Read the file content based on file type
            text_content = ""

            if file_type == FileType.PDF:
                with open(file_path, "rb") as f:
                    pdf = PdfReader(f)
                    for page in pdf.pages:
                        text_content += page.extract_text() + "\n\n"

            elif file_type == FileType.DOCX:
                doc = DocxDocument(file_path)
                for para in doc.paragraphs:
                    text_content += para.text + "\n"

            elif file_type == FileType.XLSX:
                wb = load_workbook(file_path)
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    text_content += f"Sheet: {sheet}\n"
                    for row in ws.iter_rows(values_only=True):
                        text_content += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
                    text_content += "\n"

            elif file_type == FileType.PPTX:
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content += shape.text + "\n"
                    text_content += "\n"

            else:  # TXT or other text files
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    text_content = f.read()

            # Create a LlamaIndex Document
            document = LlamaDocument(
                text=text_content,
                metadata={
                    "file_path": file_path,
                    "file_type": file_type.value,
                    "file_name": os.path.basename(file_path)
                }
            )

            return [document]

        except Exception as e:
            logger.error(f"Error loading document {file_path}: {str(e)}")
            raise

    def process_document(self, file_path: str, file_id: str, user_id: str, storage_type: str = "local") -> Dict[str, Any]:
        """
        Process a document: load, chunk, and index it.

        Args:
            file_path: Path to the file (local path or S3 URL)
            file_id: Unique ID for the file
            user_id: ID of the user who uploaded the file
            storage_type: Type of storage ("local" or "s3")

        Returns:
            Dictionary with processing results
        """
        try:
            # Handle S3 storage
            temp_file_path = None
            if storage_type == "s3":
                from app.utils.s3_storage import s3_storage
                # Extract the S3 key from the URL
                s3_key = file_path.split(f"{s3_storage.bucket_name}/")[1]
                # Download to a temporary file
                temp_file_path = os.path.join(settings.UPLOAD_DIR, f"temp_{file_id}")
                os.makedirs(os.path.dirname(temp_file_path), exist_ok=True)
                content = s3_storage.download_file(s3_key)
                with open(temp_file_path, "wb") as f:
                    f.write(content)
                # Use the temporary file for processing
                file_path = temp_file_path

            # Detect file type
            file_type = self.detect_file_type(file_path)

            # Load document
            documents = self.load_document(file_path, file_type)

            if not documents:
                raise ValueError(f"No content could be extracted from {file_path}")

            # Determine chunking strategy
            chunking_strategy = self.get_chunking_strategy(file_type)

            # Process documents based on chunking strategy
            nodes = self._chunk_documents(documents, chunking_strategy)

            # Add additional metadata
            for node in nodes:
                node.metadata.update({
                    "file_id": file_id,
                    "user_id": user_id
                })

            # Store document nodes
            if self.use_weaviate and self.weaviate_client:
                # Index documents in Weaviate
                try:
                    # Create a user-specific collection name
                    # Weaviate doesn't allow underscores in class names, so we'll replace them with hyphens
                    # Also, we'll use a shorter version of the user_id to avoid exceeding length limits
                    short_user_id = user_id.replace("-", "")[:8]
                    index_name = f"{settings.LLAMAINDEX_INDEX_NAME}{short_user_id}"
                    vector_store = WeaviateVectorStore(
                        weaviate_client=self.weaviate_client,
                        index_name=index_name,
                        text_key="text",
                        metadata_keys=["file_id", "file_type", "file_name", "user_id"]
                    )

                    # Create storage context and index
                    storage_context = StorageContext.from_defaults(vector_store=vector_store)

                    # Process nodes in batches to avoid timeouts
                    self._store_nodes_in_batches(nodes, vector_store, index_name)

                    # Return processing results
                    return {
                        "file_id": file_id,
                        "file_type": file_type.value,
                        "file_name": os.path.basename(file_path),
                        "user_id": user_id,
                        "num_chunks": len(nodes),
                        "index_name": index_name,
                        "chunking_strategy": chunking_strategy.value,
                        "status": "success",
                        "storage": "weaviate"
                    }
                except Exception as e:
                    logger.error(f"Error storing in Weaviate: {str(e)}")
                    logger.info("Falling back to in-memory storage")

            # If Weaviate is not configured or failed, store nodes in memory
            # In a real implementation, you would store these in a database
            # For now, we'll just return them
            result = {
                "file_id": file_id,
                "file_type": file_type.value,
                "file_name": os.path.basename(file_path),
                "user_id": user_id,
                "num_chunks": len(nodes),
                "nodes": nodes,
                "chunking_strategy": chunking_strategy.value,
                "status": "success",
                "storage": "memory"
            }

            # Clean up temporary file if it was created
            if temp_file_path and os.path.exists(temp_file_path):
                try:
                    os.remove(temp_file_path)
                    logger.info(f"Removed temporary file {temp_file_path}")
                except Exception as e:
                    logger.warning(f"Failed to remove temporary file {temp_file_path}: {str(e)}")

            return result

        except Exception as e:
            logger.error(f"Error processing document {file_path}: {str(e)}")
            # Clean up temporary file if it was created
            if 'temp_file_path' in locals() and temp_file_path and os.path.exists(temp_file_path):
                try:
                    os.remove(temp_file_path)
                    logger.info(f"Removed temporary file {temp_file_path}")
                except Exception as cleanup_error:
                    logger.warning(f"Failed to remove temporary file {temp_file_path}: {str(cleanup_error)}")

            return {
                "file_id": file_id,
                "file_path": file_path,
                "user_id": user_id,
                "status": "error",
                "error": str(e)
            }

    def _chunk_documents(self, documents: List[LlamaDocument], strategy: ChunkingStrategy) -> List[Any]:
        """
        Chunk documents based on the specified strategy.

        Args:
            documents: List of LlamaIndex Document objects
            strategy: ChunkingStrategy enum value

        Returns:
            List of document nodes
        """
        if strategy == ChunkingStrategy.FIXED_SIZE:
            return self.fixed_size_parser.get_nodes_from_documents(documents)

        elif strategy == ChunkingStrategy.TOPIC_BASED:
            return self.topic_based_parser.get_nodes_from_documents(documents)

        elif strategy == ChunkingStrategy.HYBRID:
            # For hybrid strategy, we'll use a combination of approaches
            # First, try topic-based chunking
            topic_nodes = self.topic_based_parser.get_nodes_from_documents(documents)

            # Check if the chunks are reasonable sizes
            large_chunks = [node for node in topic_nodes if len(node.text) > settings.MAX_CHUNK_SIZE]
            small_chunks = [node for node in topic_nodes if len(node.text) < settings.MIN_CHUNK_SIZE]

            # If we have too many large or small chunks, fall back to fixed-size chunking
            if len(large_chunks) > len(topic_nodes) * 0.3 or len(small_chunks) > len(topic_nodes) * 0.5:
                return self.fixed_size_parser.get_nodes_from_documents(documents)

            return topic_nodes

        else:
            # Default to fixed-size chunking
            return self.fixed_size_parser.get_nodes_from_documents(documents)

    def _store_nodes_in_batches(self, nodes: List, vector_store, index_name: str) -> None:
        """
        Store nodes in Weaviate using batched processing to avoid timeouts.

        Args:
            nodes: List of nodes to store
            vector_store: Vector store instance
            index_name: Name of the index/collection
        """
        try:
            # Get the total number of nodes
            total_nodes = len(nodes)
            logger.info(f"Starting batch processing of {total_nodes} nodes to collection {index_name}")

            # Calculate number of batches
            batch_size = settings.WEAVIATE_BATCH_SIZE
            num_batches = (total_nodes + batch_size - 1) // batch_size  # Ceiling division

            # Process nodes in batches
            for batch_idx in range(num_batches):
                start_idx = batch_idx * batch_size
                end_idx = min(start_idx + batch_size, total_nodes)
                batch_nodes = nodes[start_idx:end_idx]

                logger.info(f"Processing batch {batch_idx + 1}/{num_batches} with {len(batch_nodes)} nodes")

                # Create a temporary storage context for this batch
                batch_storage_context = StorageContext.from_defaults(
                    vector_store=vector_store
                )

                # Process the batch with retries
                retry_count = 0
                max_retries = settings.WEAVIATE_MAX_RETRIES
                success = False

                while retry_count < max_retries and not success:
                    try:
                        # Create a temporary index for this batch
                        VectorStoreIndex(
                            nodes=batch_nodes,
                            storage_context=batch_storage_context,
                        )
                        success = True
                        logger.info(f"Successfully processed batch {batch_idx + 1}/{num_batches}")
                    except Exception as e:
                        retry_count += 1
                        logger.warning(f"Batch {batch_idx + 1} attempt {retry_count} failed: {str(e)}")
                        if retry_count < max_retries:
                            logger.info(f"Retrying batch {batch_idx + 1} (attempt {retry_count + 1}/{max_retries})...")
                            # Wait before retrying with exponential backoff
                            time.sleep(2 ** retry_count)
                        else:
                            logger.error(f"Failed to process batch {batch_idx + 1} after {max_retries} attempts")
                            # Continue with next batch instead of failing the entire process

            logger.info(f"Completed batch processing of {total_nodes} nodes to collection {index_name}")
        except Exception as e:
            logger.error(f"Error in batch processing: {str(e)}")
            # Don't raise the exception to allow the process to continue
            # The document will be marked as processed even if some batches failed

# Create a singleton instance
document_processor = DocumentProcessor()
